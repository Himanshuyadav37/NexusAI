import os
import zipfile
import base64
import logging
import re
import urllib.parse
from pathlib import Path
from typing import List, Dict, Any
import requests
from config import settings

# Try imports, log warnings if somehow they fail to load
try:
    import pypdf
except ImportError:
    pypdf = None

try:
    import docx
except ImportError:
    docx = None

try:
    import pptx
except ImportError:
    pptx = None

try:
    import openpyxl
except ImportError:
    openpyxl = None

logger = logging.getLogger(__name__)

# ==========================================
# Gemini OCR & Multimodal Extractors
# ==========================================
def extract_text_via_gemini(file_bytes: bytes, mime_type: str) -> str:
    """Send document/image directly to Gemini API to extract text content."""
    gemini_key = getattr(settings, "GEMINI_API_KEY", "")
    if not gemini_key:
        logger.warning("GEMINI_API_KEY is not set. Cannot perform Gemini OCR.")
        return "[Error: GEMINI_API_KEY missing for OCR]"
        
    url = f"https://generativelanguage.googleapis.com/v1beta/models/gemini-1.5-flash:generateContent?key={gemini_key}"
    base64_data = base64.b64encode(file_bytes).decode("utf-8")
    
    payload = {
        "contents": [
            {
                "parts": [
                    {"text": "Extract all text from this document or image. Keep the layout structure as much as possible. Return ONLY the raw extracted text. Do not add any conversational introductions, summaries, or formatting notes."},
                    {
                        "inline_data": {
                            "mime_type": mime_type,
                            "data": base64_data
                        }
                    }
                ]
            }
        ],
        "generationConfig": {
            "temperature": 0.1
        }
    }
    
    try:
        response = requests.post(url, json=payload, timeout=45)
        response.raise_for_status()
        res_json = response.json()
        text = res_json["contents"][0]["parts"][0]["text"]
        return text.strip()
    except Exception as e:
        logger.error(f"Gemini OCR text extraction failed: {e}")
        return f"[OCR Extraction Failed: {str(e)}]"

# ==========================================
# Format-Specific Document Parsers
# ==========================================
def parse_pdf(file_path: Path) -> List[Dict[str, Any]]:
    pages_data = []
    filename = file_path.name
    
    if pypdf is None:
        logger.error("pypdf is not installed. PDF parsing unavailable.")
        return [{"text": "[Error: pypdf not installed]", "page_num": 1, "filename": filename}]

    try:
        with open(file_path, "rb") as f:
            reader = pypdf.PdfReader(f)
            total_pages = len(reader.pages)
            
            # Check if pdf is scanned (has very little text in first few pages)
            check_text = ""
            for i in range(min(3, total_pages)):
                check_text += reader.pages[i].extract_text() or ""
                
            if len(check_text.strip()) < 100:
                logger.info(f"PDF {filename} appears to be scanned (extracted text length={len(check_text)}). Using Gemini OCR.")
                f.seek(0)
                pdf_bytes = f.read()
                ocr_text = extract_text_via_gemini(pdf_bytes, "application/pdf")
                # Split OCR text roughly into single page entries or store as one page
                return [{"text": ocr_text, "page_num": 1, "filename": filename}]
            
            # Extract page by page normally
            for idx, page in enumerate(reader.pages):
                page_num = idx + 1
                text = page.extract_text() or ""
                
                # If a specific page is blank/scanned, we can run OCR on just that page
                if not text.strip():
                    logger.info(f"Page {page_num} of {filename} is blank. Attempting page-specific OCR...")
                    try:
                        # Extract single page as a PDF bytes representation
                        writer = pypdf.PdfWriter()
                        writer.add_page(page)
                        from io import BytesIO
                        page_io = BytesIO()
                        writer.write(page_io)
                        page_bytes = page_io.getvalue()
                        text = extract_text_via_gemini(page_bytes, "application/pdf")
                    except Exception as page_err:
                        logger.warning(f"Failed page-specific OCR for page {page_num}: {page_err}")
                
                pages_data.append({
                    "text": text.strip(),
                    "page_num": page_num,
                    "filename": filename
                })
    except Exception as e:
        logger.error(f"Error parsing PDF {file_path}: {e}")
        # Fallback to direct OCR on the whole file
        try:
            with open(file_path, "rb") as f:
                pdf_bytes = f.read()
            ocr_text = extract_text_via_gemini(pdf_bytes, "application/pdf")
            return [{"text": ocr_text, "page_num": 1, "filename": filename}]
        except Exception as fallback_err:
            logger.error(f"Fallback OCR also failed: {fallback_err}")
            pages_data.append({
                "text": f"[Failed to parse PDF: {str(e)}]",
                "page_num": 1,
                "filename": filename
            })
            
    return pages_data

def parse_docx(file_path: Path) -> List[Dict[str, Any]]:
    filename = file_path.name
    if docx is None:
        return [{"text": "[Error: python-docx not installed]", "page_num": 1, "filename": filename}]
    try:
        doc = docx.Document(file_path)
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        for table in doc.tables:
            for row in table.rows:
                row_text = [cell.text for cell in row.cells]
                full_text.append(" | ".join(row_text))
        return [{"text": "\n".join(full_text), "page_num": 1, "filename": filename}]
    except Exception as e:
        logger.error(f"Error parsing DOCX {file_path}: {e}")
        return [{"text": f"[Error parsing DOCX: {e}]", "page_num": 1, "filename": filename}]

def parse_pptx(file_path: Path) -> List[Dict[str, Any]]:
    filename = file_path.name
    if pptx is None:
        return [{"text": "[Error: python-pptx not installed]", "page_num": 1, "filename": filename}]
    try:
        presentation = pptx.Presentation(file_path)
        slides_data = []
        for idx, slide in enumerate(presentation.slides):
            slide_text = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    slide_text.append(shape.text)
            slides_data.append({
                "text": "\n".join(slide_text),
                "page_num": idx + 1,
                "filename": filename
            })
        return slides_data
    except Exception as e:
        logger.error(f"Error parsing PPTX {file_path}: {e}")
        return [{"text": f"[Error parsing PPTX: {e}]", "page_num": 1, "filename": filename}]

def parse_xlsx(file_path: Path) -> List[Dict[str, Any]]:
    filename = file_path.name
    if openpyxl is None:
        return [{"text": "[Error: openpyxl not installed]", "page_num": 1, "filename": filename}]
    try:
        wb = openpyxl.load_workbook(file_path, read_only=True, data_only=True)
        sheets_text = []
        for sheet_name in wb.sheetnames:
            sheet = wb[sheet_name]
            sheets_text.append(f"### Sheet: {sheet_name}")
            for row in sheet.iter_rows(values_only=True):
                # Filter out completely empty rows
                if any(cell is not None for cell in row):
                    row_str = " | ".join(str(cell) if cell is not None else "" for cell in row)
                    sheets_text.append(row_str)
        return [{"text": "\n".join(sheets_text), "page_num": 1, "filename": filename}]
    except Exception as e:
        logger.error(f"Error parsing XLSX {file_path}: {e}")
        return [{"text": f"[Error parsing XLSX: {e}]", "page_num": 1, "filename": filename}]

def parse_text(file_path: Path) -> List[Dict[str, Any]]:
    filename = file_path.name
    try:
        with open(file_path, "r", encoding="utf-8", errors="ignore") as f:
            content = f.read()
        return [{"text": content, "page_num": 1, "filename": filename}]
    except Exception as e:
        logger.error(f"Error parsing text file {file_path}: {e}")
        return [{"text": f"[Error parsing file: {e}]", "page_num": 1, "filename": filename}]

def parse_image(file_path: Path) -> List[Dict[str, Any]]:
    filename = file_path.name
    suffix = file_path.suffix.lower()
    mime_map = {
        ".png": "image/png",
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".webp": "image/webp"
    }
    mime_type = mime_map.get(suffix, "image/png")
    try:
        with open(file_path, "rb") as f:
            image_bytes = f.read()
        ocr_text = extract_text_via_gemini(image_bytes, mime_type)
        return [{"text": ocr_text, "page_num": 1, "filename": filename}]
    except Exception as e:
        logger.error(f"Error parsing image {file_path}: {e}")
        return [{"text": f"[Image OCR Failed: {e}]", "page_num": 1, "filename": filename}]

# ==========================================
# Web Crawling & GitHub Repository Downloader
# ==========================================
def parse_url(url: str) -> List[Dict[str, Any]]:
    """Fetch website URL, strip tags, and structure raw text."""
    try:
        headers = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) NexusAIRAG/1.0"}
        response = requests.get(url, headers=headers, timeout=15)
        response.raise_for_status()
        
        # Clean HTML (simple regex-based tag stripper + spacing)
        html = response.text
        # Remove script and style elements
        html = re.sub(r"<(script|style).*?>.*?</\1>", "", html, flags=re.DOTALL | re.IGNORECASE)
        # Remove tags
        text = re.sub(r"<[^>]*>", " ", html)
        # Standardize whitespaces
        text = re.sub(r"\s+", " ", text).strip()
        
        parsed_url = urllib.parse.urlparse(url)
        filename = parsed_url.netloc + parsed_url.path
        if filename.endswith("/"):
            filename = filename[:-1]
            
        return [{"text": f"Source URL: {url}\n\n{text}", "page_num": 1, "filename": filename}]
    except Exception as e:
        logger.error(f"Error fetching/parsing URL {url}: {e}")
        return [{"text": f"[Error fetching URL {url}: {e}]", "page_num": 1, "filename": "url_error"}]

def parse_github(repo_url: str) -> List[Dict[str, Any]]:
    """Download public GitHub repo zip, parse text, and clean up."""
    # Convert github.com/user/repo to zipball
    match = re.search(r"github\.com/([^/]+)/([^/]+)", repo_url)
    if not match:
        return [{"text": f"[Invalid GitHub repo URL: {repo_url}]", "page_num": 1, "filename": "github_error"}]
        
    owner, repo = match.group(1), match.group(2)
    # Remove trailing .git if present
    if repo.endswith(".git"):
        repo = repo[:-4]
        
    zip_url = f"https://github.com/{owner}/{repo}/zipball/main"
    
    # Create temporary file inside project workspace folder
    temp_zip = Path(__file__).resolve().parents[2] / f"temp_git_{owner}_{repo}.zip"
    try:
        logger.info(f"Downloading GitHub ZIP from {zip_url}...")
        headers = {}
        github_token = getattr(settings, "GITHUB_TOKEN", "")
        if github_token:
            headers["Authorization"] = f"token {github_token}"
            
        response = requests.get(zip_url, headers=headers, timeout=60)
        # Fallback to master if main fails
        if response.status_code != 200:
            logger.info("Main branch download failed, trying master branch...")
            zip_url = f"https://github.com/{owner}/{repo}/zipball/master"
            response = requests.get(zip_url, headers=headers, timeout=60)
            
        response.raise_for_status()
        with open(temp_zip, "wb") as f:
            f.write(response.content)
            
        # Parse zip recursively
        pages = parse_zip(temp_zip)
        return pages
    except Exception as e:
        logger.error(f"GitHub repository download failed: {e}")
        return [{"text": f"[Error importing GitHub repository {repo_url}: {e}]", "page_num": 1, "filename": "github_error"}]
    finally:
        if temp_zip.exists():
            try:
                os.remove(temp_zip)
            except Exception:
                pass

# ==========================================
# Main File Router & Ingestion Engine
# ==========================================
def parse_file(file_path: Path) -> List[Dict[str, Any]]:
    """Route document parsing based on extension."""
    suffix = file_path.suffix.lower()
    
    if suffix == ".pdf":
        return parse_pdf(file_path)
    elif suffix == ".docx":
        return parse_docx(file_path)
    elif suffix == ".pptx":
        return parse_pptx(file_path)
    elif suffix in [".xlsx", ".xls"]:
        return parse_xlsx(file_path)
    elif suffix in [".png", ".jpg", ".jpeg", ".webp"]:
        return parse_image(file_path)
    elif suffix in [".zip"]:
        return parse_zip(file_path)
    else:
        # Fallback to raw text reader for md, txt, csv, code source files, etc.
        return parse_text(file_path)

def parse_zip(zip_path: Path) -> List[Dict[str, Any]]:
    """Unzip dynamically inside workspace and index recursively."""
    extracted_docs = []
    # Temporary extract dir in the workspace
    extract_dir = zip_path.parent / f"extracted_{zip_path.stem}"
    extract_dir.mkdir(parents=True, exist_ok=True)
    
    try:
        with zipfile.ZipFile(zip_path, "r") as zip_ref:
            zip_ref.extractall(extract_dir)
            
        # Recursively walk the folder
        extracted_docs = parse_directory(extract_dir)
    except Exception as e:
        logger.error(f"Error unzipping {zip_path}: {e}")
        extracted_docs.append({
            "text": f"[Error unzipping archive: {e}]",
            "page_num": 1,
            "filename": zip_path.name
        })
    finally:
        # Clean up files
        import shutil
        if extract_dir.exists():
            try:
                shutil.rmtree(extract_dir)
            except Exception as clean_err:
                logger.warning(f"Failed to clean extracted ZIP dir: {clean_err}")
                
    return extracted_docs

def parse_directory(dir_path: Path) -> List[Dict[str, Any]]:
    """Recursively traverse directory and parse files."""
    all_pages = []
    # Skip blacklisted directories (standard code repo boilerplate)
    skip_dirs = {".git", "node_modules", "__pycache__", "venv", "dist", "build", ".idea"}
    skip_extensions = {".zip", ".tar", ".gz", ".rar", ".exe", ".dll", ".so", ".pyc", ".png", ".jpg", ".jpeg", ".webp"} # skip images in recursive repository scanning to prevent spamming Gemini OCR calls unless specifically uploaded
    
    for root, dirs, files in os.walk(dir_path):
        # In-place modify dirs to skip walking blacklisted folders
        dirs[:] = [d for d in dirs if d not in skip_dirs]
        
        for file in files:
            file_path = Path(root) / file
            if file_path.suffix.lower() in skip_extensions:
                continue
                
            # Compute relative path in ZIP/Folder for citation naming
            try:
                rel_path = file_path.relative_to(dir_path)
                display_name = str(rel_path)
            except ValueError:
                display_name = file_path.name
                
            file_pages = parse_file(file_path)
            for p in file_pages:
                # Update display name to retain zip folder structures
                p["filename"] = display_name
                all_pages.append(p)
                
    return all_pages
